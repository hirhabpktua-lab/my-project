from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def add_slide_number(slide, num, total):
    """右下角页码，属于细节点缀"""
    left = Inches(8.5)
    top = Inches(7.1)
    txBox = slide.shapes.add_textbox(left, top, Inches(1.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.RIGHT


def add_accent_bar(slide, top=Inches(0.05)):
    """顶部装饰线，让页面不那么素"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, Inches(10), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x2B, 0x57, 0x9A)
    shape.line.fill.background()


def set_slide_bg(slide, color_hex="F5F5F5"):
    """统一浅灰背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4)))


prs = Presentation()
prs.slide_width = Inches(10)    # 16:9 宽屏
prs.slide_height = Inches(7.5)
TOTAL = 6


# ── Slide 1: 封面 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
set_slide_bg(slide, "2B579A")

# 大标题
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1.2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "New Project"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

# 副标题
txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.3), Inches(7), Inches(0.8))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "项目演示与功能概览"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(0xCC, 0xD5, 0xE8)
p2.alignment = PP_ALIGN.CENTER

# 底部信息条
bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(10), Inches(0.03)
)
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
bar.line.fill.background()

txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6.7), Inches(8), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "2026.05 · GitHub: hirhabpktua-lab/my-project"
p3.font.size = Pt(11)
p3.font.color.rgb = RGBColor(0x99, 0xAA, 0xCC)
p3.alignment = PP_ALIGN.CENTER

add_slide_number(slide, 1, TOTAL)


# ── Slide 2: 目录 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(5), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "目录"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2B, 0x57, 0x9A)

items = [
    ("01", "项目背景与目标"),
    ("02", "技术架构一览"),
    ("03", "核心功能演示"),
    ("04", "下一步计划"),
]
for i, (num, title) in enumerate(items):
    y = Inches(1.9) + Inches(i * 1.15)
    # 编号圆圈
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.2), y, Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x2B, 0x57, 0x9A)
    circle.line.fill.background()
    ctf = circle.text_frame
    cp = ctf.paragraphs[0]
    cp.text = num
    cp.font.size = Pt(14)
    cp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cp.font.bold = True
    cp.alignment = PP_ALIGN.CENTER

    # 标题文字
    tbox = slide.shapes.add_textbox(Inches(2.0), y + Inches(0.05), Inches(6), Inches(0.45))
    tf_item = tbox.text_frame
    ip = tf_item.paragraphs[0]
    ip.text = title
    ip.font.size = Pt(20)
    ip.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

add_slide_number(slide, 2, TOTAL)


# ── Slide 3: 项目背景 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "项目背景与目标"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2B, 0x57, 0x9A)

content = [
    ("背景", "本项目旨在构建一套高效的开发工作流，集成代码审查、自动化测试、前端设计等核心能力，提升团队开发效率与代码质量。"),
    ("目标", "打造一站式开发工具链，让每次提交都经过自动化审查、每次部署都有质量保证，实现真正的\"放心推\"。"),
]
for i, (label, text) in enumerate(content):
    y = Inches(1.9) + Inches(i * 2.0)

    # 标签色块
    lbl = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(1.3), Inches(0.45)
    )
    lbl.fill.solid()
    lbl.fill.fore_color.rgb = RGBColor(0x2B, 0x57, 0x9A)
    lbl.line.fill.background()
    ltf = lbl.text_frame
    lp = ltf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(14)
    lp.font.bold = True
    lp.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    lp.alignment = PP_ALIGN.CENTER

    # 内容文本
    tbox = slide.shapes.add_textbox(Inches(2.4), y + Inches(0.05), Inches(6.5), Inches(1.2))
    tf_item = tbox.text_frame
    tf_item.word_wrap = True
    ip = tf_item.paragraphs[0]
    ip.text = text
    ip.font.size = Pt(16)
    ip.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
    ip.space_after = Pt(6)

add_slide_number(slide, 3, TOTAL)


# ── Slide 4: 技术架构 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "技术架构一览"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2B, 0x57, 0x9A)

categories = [
    ("语言 & 运行时", "Node.js v24 · Python 3.12 · TypeScript"),
    ("前端", "React · Next.js · Playwright 浏览器自动化"),
    ("后端 & 数据", "Express · PostgreSQL · Redis"),
    ("开发工具", "Git · ESLint · Prettier · pytest · black · ruff"),
    ("AI / LLM", "Claude Agent SDK · Anthropic API · Context7 文档查询"),
]
for i, (cat, techs) in enumerate(categories):
    y = Inches(1.7) + Inches(i * 0.95)

    cat_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(2.2), Inches(0.4)
    )
    cat_box.fill.solid()
    cat_box.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)
    cat_box.line.fill.background()
    tfc = cat_box.text_frame
    ccp = tfc.paragraphs[0]
    ccp.text = cat
    ccp.font.size = Pt(13)
    ccp.font.bold = True
    ccp.font.color.rgb = RGBColor(0x2B, 0x57, 0x9A)
    ccp.alignment = PP_ALIGN.CENTER

    tbox = slide.shapes.add_textbox(Inches(3.3), y + Inches(0.02), Inches(6), Inches(0.4))
    tfi = tbox.text_frame
    ipp = tfi.paragraphs[0]
    ipp.text = techs
    ipp.font.size = Pt(14)
    ipp.font.color.rgb = RGBColor(0x44, 0x44, 0x44)

add_slide_number(slide, 4, TOTAL)


# ── Slide 5: 核心功能 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_accent_bar(slide)

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8), Inches(0.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "核心功能演示"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2B, 0x57, 0x9A)

features = [
    ("🧠", "智能代码审查", "Git 提交自动触发审查，覆盖安全漏洞、代码风格、逻辑错误"),
    ("🔄", "自动化工作流", "Git Hooks + CI/CD 全链路，从编码到部署零人工干预"),
    ("🎨", "前端设计系统", "内置 Playwright 截图对比，CSS 回归一跑即知"),
    ("📦", "PR 审查工具包", "类型设计分析、静默失败检测、测试覆盖审查，一键出报告"),
]
for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 4.0)
    y = Inches(1.8) + Inches(row * 2.4)

    # 卡片背景
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.0)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(0.5)
    # 轻微阴影效果 — 通过偏移深色矩形在后面模拟
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.03), y + Inches(0.03), Inches(3.7), Inches(2.0)
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shadow.line.fill.background()
    # 把阴影移到卡片后面 (z-order: 先添加的在底层，我们用顺序控制)

    # 图标
    ibox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5))
    itf = ibox.text_frame
    ip = itf.paragraphs[0]
    ip.text = icon
    ip.font.size = Pt(24)

    # 标题
    tbox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.75), Inches(3.2), Inches(0.4))
    ttf = tbox.text_frame
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(16)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0x2B, 0x57, 0x9A)

    # 描述
    dbox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(1.15), Inches(3.2), Inches(0.75))
    dtf = dbox.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(12)
    dp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

add_slide_number(slide, 5, TOTAL)


# ── Slide 6: 结语 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, "2B579A")

txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(8), Inches(1.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "谢谢观看"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(7), Inches(0.7))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "下一步：完善测试覆盖 · 接入 CI/CD · 上线生产"
p2.font.size = Pt(16)
p2.font.color.rgb = RGBColor(0xCC, 0xD5, 0xE8)
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(0.6))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "github.com/hirhabpktua-lab/my-project"
p3.font.size = Pt(13)
p3.font.color.rgb = RGBColor(0x99, 0xAA, 0xCC)
p3.alignment = PP_ALIGN.CENTER

add_slide_number(slide, 6, TOTAL)


# ── 保存 ──
output_path = r"C:\Users\Administrator\Documents\New project\New-Project-演示.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
